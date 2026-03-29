import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

def add_title_slide(title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.title
    subtitle_box = slide.placeholders[1]
    title_box.text = title
    subtitle_box.text = subtitle

def add_bullet_slide(title, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.title
    body_box = slide.placeholders[1]
    title_box.text = title
    tf = body_box.text_frame
    for i, bullet in enumerate(bullets):
        if i == 0:
            tf.text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0

add_title_slide("CORTEX Dictionary Lab", "Advanced Corpus-Driven Dictionary Editor\nFeatures Overview")

add_bullet_slide("Slide 2: Corpus Management & Indexing", [
    "Supports Monolingual and Parallel Corpora in multiple formats (XML, TXT, CSV, XLSX).",
    "Dynamic ingest and database indexing via DuckDB for fast querying.",
    "Easily choose between 'Built-In' corpora on disk or upload files securely.",
    "Handles both Source and Target alignments for Parallel corpus study."
])

add_bullet_slide("Slide 3: Comprehensive Phrase & Keyword Search", [
    "Search single words or complex multi-word phrases.",
    "Integrated Autocomplete based on loaded corpora and metadata.",
    "Constituent word breakdown & automatically linked phrases.",
    "Instant feedback on missing terms or hidden terms under current filters."
])

add_bullet_slide("Slide 4: Advanced Statistical Context (PMW & Zipf)", [
    "Frequency metrics calculated contextually per selected corpus.",
    "PMW (Per Million Words) indicator with relative bar chart UI.",
    "Zipf scale visualization (1-5) for intuitive word rarity understanding.",
    "Rank collocation strength explicitly evaluated for nodes."
])

add_bullet_slide("Slide 5: Robust Collocation & N-Grams", [
    "Identifies Top-20 Collocates based on statistical metrics (e.g., Log-Likelihood).",
    "Forward and Backward N-gram search showing the 'Phrase + Word' or 'Word + Phrase' contexts.",
    "Interactive Radial Network Graph (LancsBox style) visualizing collocates:",
    " - Distance indicates strength.",
    " - Size indicates frequency.",
    " - Left vs. right positional dominance."
])

add_bullet_slide("Slide 6: KWIC (Key Word In Context) Examples", [
    "Interactive, heavily-formatted Concordance Lines.",
    "Highlights Node word in distinct styles.",
    "Specific Collocate usage examples isolating sentences where a node and collocate appear together.",
    "Translates parallel lines automatically below sentence examples."
])

add_bullet_slide("Slide 7: Vocabulary Profiling & Badges", [
    "Intelligently evaluates vocabulary using POS tags.",
    "Integrates seamlessly with CEFR, GSL, NGSL, AWL metrics natively via python libraries.",
    "Allows custom User-Defined Wordlists (via .txt uploads) applied across languages.",
    "Dynamic rendering of visual badges on dictionary entries."
])

add_bullet_slide("Slide 8: AI-Powered Dictionary Generation", [
    "Integrate Local AI (Ollama) or Cloud AI (Google Gemini).",
    "Automatically generate missing Definitions, Pronunciations, and Collocate examples.",
    "Extract KWIC samples to feed into contextual AI prompts.",
    "Provides real-time 'Accept or Dismiss' options for AI-suggested content."
])

add_bullet_slide("Slide 9: Corpus Statistic & Frequency Tab", [
    "Global dashboard displaying Total Tokens, Lemmas, and Unique POS Tags.",
    "Displays a fully comprehensive, downloadable Excel (xlsx) Frequency List.",
    "Interactive Pandas dataframe display inside the view with clickable navigation.",
    "Easily inspect how widespread a term is across the dataset."
])

add_bullet_slide("Slide 10: Personal Overrides & Data Persistence", [
    "Create and edit new lexical dictionary sense fields dynamically via the UI.",
    "Save changes automatically to a personalized local json configuration.",
    "Changes made persistently reflect during search and across Corpus Statistics.",
    "Allows for complete collaborative editing workflows within a lab environment."
])

prs.save('CORTEX_Dictionary_Lab_Features.pptx')
print("Presentation created successfully as 'CORTEX_Dictionary_Lab_Features.pptx'")
